"""Tests for document text extraction utilities."""

from pathlib import Path
from zipfile import ZipFile

import pytest

from nanobot.utils.document import (
    SUPPORTED_EXTENSIONS,
    PdfSafetyError,
    _is_text_extension,
    extract_pdf_pages,
    extract_text,
)


class TestSupportedExtensions:
    """Test the SUPPORTED_EXTENSIONS constant."""

    def test_supported_extensions_include_common_formats(self):
        """Test that common document formats are included."""
        # Document formats
        assert ".pdf" in SUPPORTED_EXTENSIONS
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".xlsx" in SUPPORTED_EXTENSIONS
        assert ".pptx" in SUPPORTED_EXTENSIONS

        # Text formats
        assert ".txt" in SUPPORTED_EXTENSIONS
        assert ".md" in SUPPORTED_EXTENSIONS
        assert ".csv" in SUPPORTED_EXTENSIONS
        assert ".json" in SUPPORTED_EXTENSIONS
        assert ".yaml" in SUPPORTED_EXTENSIONS
        assert ".yml" in SUPPORTED_EXTENSIONS

        # Image formats
        assert ".png" in SUPPORTED_EXTENSIONS
        assert ".jpg" in SUPPORTED_EXTENSIONS
        assert ".jpeg" in SUPPORTED_EXTENSIONS


class TestExtractText:
    """Test the extract_text function."""

    def test_extract_text_unsupported_returns_none(self, tmp_path: Path):
        """Test that unsupported file types return None."""
        unsupported_file = tmp_path / "file.xyz"
        unsupported_file.write_text("content")

        result = extract_text(unsupported_file)
        assert result is None

    def test_extract_text_file_not_found(self, tmp_path: Path):
        """Test that non-existent files return error string."""
        missing_file = tmp_path / "nonexistent.txt"

        result = extract_text(missing_file)
        assert result is not None
        assert "[error: file not found:" in result

    def test_extract_text_txt_file(self, tmp_path: Path):
        """Test extracting text from a .txt file."""
        txt_file = tmp_path / "test.txt"
        content = "Hello, world!\nThis is a test."
        txt_file.write_text(content, encoding="utf-8")

        result = extract_text(txt_file)
        assert result == content

    def test_extract_text_txt_file_with_truncation(self, tmp_path: Path):
        """Test that large text files are truncated."""
        txt_file = tmp_path / "large.txt"
        # Create content larger than _MAX_TEXT_LENGTH
        content = "x" * 300_000
        txt_file.write_text(content, encoding="utf-8")

        result = extract_text(txt_file)
        assert len(result) < 300_000
        assert "(truncated," in result
        assert "chars total)" in result

    def test_extract_text_md_file(self, tmp_path: Path):
        """Test extracting text from a .md file."""
        md_file = tmp_path / "test.md"
        content = "# Header\n\nSome markdown content."
        md_file.write_text(content, encoding="utf-8")

        result = extract_text(md_file)
        assert result == content

    def test_extract_text_csv_file(self, tmp_path: Path):
        """Test extracting text from a .csv file."""
        csv_file = tmp_path / "test.csv"
        content = "name,age\nAlice,30\nBob,25"
        csv_file.write_text(content, encoding="utf-8")

        result = extract_text(csv_file)
        assert result == content

    def test_extract_text_json_file(self, tmp_path: Path):
        """Test extracting text from a .json file."""
        json_file = tmp_path / "test.json"
        content = '{"key": "value", "number": 42}'
        json_file.write_text(content, encoding="utf-8")

        result = extract_text(json_file)
        assert result == content

    def test_extract_text_xlsx(self, tmp_path: Path):
        """Test extracting text from an .xlsx file."""
        from openpyxl import Workbook

        xlsx_file = tmp_path / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["A2"] = "Alice"
        ws["B2"] = 30
        ws["A3"] = "Bob"
        ws["B3"] = 25

        # Add a second sheet
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Product"
        ws2["B1"] = "Price"
        ws2["A2"] = "Widget"
        ws2["B2"] = 9.99

        wb.save(xlsx_file)
        wb.close()

        result = extract_text(xlsx_file)
        assert result is not None
        assert "--- Sheet: Sheet1 ---" in result
        assert "--- Sheet: Sheet2 ---" in result
        assert "Alice" in result
        assert "Bob" in result
        assert "Widget" in result
        assert "9.99" in result

    def test_extract_text_xlsx_empty_sheet(self, tmp_path: Path):
        """Test extracting text from an .xlsx file with empty sheets."""
        from openpyxl import Workbook

        xlsx_file = tmp_path / "empty.xlsx"
        wb = Workbook()
        # Clear the default sheet
        wb.remove(wb.active)
        # Add an empty sheet
        wb.create_sheet("EmptySheet")
        wb.save(xlsx_file)
        wb.close()

        result = extract_text(xlsx_file)
        # Empty sheets should return empty string or header only
        assert result == "--- Sheet: EmptySheet ---" or result == ""

    def test_extract_text_docx(self, tmp_path: Path):
        """Test extracting text from a .docx file."""
        from docx import Document

        docx_file = tmp_path / "test.docx"
        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is paragraph one.")
        doc.add_paragraph("This is paragraph two.")
        doc.save(docx_file)

        result = extract_text(docx_file)
        assert result is not None
        assert "Test Document" in result
        assert "This is paragraph one." in result
        assert "This is paragraph two." in result

    def test_extract_text_docx_preserves_paragraph_and_table_order(self, tmp_path: Path):
        """DOCX forms commonly keep nearly all meaningful content in tables."""
        from docx import Document

        docx_file = tmp_path / "form.docx"
        doc = Document()
        doc.add_paragraph("Applicant details")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Ada Lovelace"
        table.cell(1, 0).text = "Project"
        table.cell(1, 1).text = "Analytical Engine"
        doc.add_paragraph("End of form")
        doc.save(docx_file)

        result = extract_text(docx_file)

        assert result is not None
        assert "Name\tAda Lovelace" in result
        assert "Project\tAnalytical Engine" in result
        assert result.index("Applicant details") < result.index("Name\tAda Lovelace")
        assert result.index("Analytical Engine") < result.index("End of form")

    def test_extract_text_docx_preserves_nested_table_text(self, tmp_path: Path):
        """Nested layout tables must not silently drop form fields."""
        from docx import Document

        docx_file = tmp_path / "nested-form.docx"
        doc = Document()
        outer_cell = doc.add_table(rows=1, cols=1).cell(0, 0)
        outer_cell.add_paragraph("Contact")
        nested = outer_cell.add_table(rows=1, cols=2)
        nested.cell(0, 0).text = "Email"
        nested.cell(0, 1).text = "ada@example.com"
        doc.save(docx_file)

        result = extract_text(docx_file)

        assert result is not None
        assert "Contact" in result
        assert "Email" in result
        assert "ada@example.com" in result
        assert result.index("Contact") < result.index("Email") < result.index("ada@example.com")

    def test_extract_text_docx_does_not_expand_grid_spans(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        """Physical cells avoid python-docx's eager gridSpan expansion."""
        from docx import Document
        from docx.table import _Row

        docx_file = tmp_path / "merged.docx"
        doc = Document()
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).merge(table.cell(0, 1)).text = "Only once"
        doc.save(docx_file)

        def fail_on_expansion(_row: _Row):
            pytest.fail("row.cells expands gridSpan before extraction can apply a bound")

        monkeypatch.setattr(_Row, "cells", property(fail_on_expansion))

        assert extract_text(docx_file) == "Only once"

    def test_extract_text_docx_keeps_vertical_merges_compact(self, tmp_path: Path):
        """Vertically merged labels appear once without shifting later columns."""
        from docx import Document

        docx_file = tmp_path / "vertical-merge.docx"
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).merge(table.cell(1, 0)).text = "Group"
        table.cell(0, 1).text = "First"
        table.cell(1, 1).text = "Second"
        doc.save(docx_file)

        assert extract_text(docx_file) == "Group\tFirst\n\tSecond"

    def test_extract_text_docx_bounds_physical_table_cells(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        """Large tables fail safely even when their text output would be empty."""
        from docx import Document

        from nanobot.utils import document as document_utils

        docx_file = tmp_path / "too-many-cells.docx"
        doc = Document()
        doc.add_table(rows=1, cols=2)
        doc.save(docx_file)
        monkeypatch.setattr(document_utils, "_MAX_DOCX_TABLE_CELLS", 1)

        result = extract_text(docx_file)

        assert result is not None
        assert result.startswith("[error: unsafe DOCX:")

    def test_extract_text_docx_bounds_table_nesting(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        """Deeply nested tables fail safely instead of recursing without a bound."""
        from docx import Document

        from nanobot.utils import document as document_utils

        docx_file = tmp_path / "nested-too-deep.docx"
        doc = Document()
        outer_cell = doc.add_table(rows=1, cols=1).cell(0, 0)
        outer_cell.add_table(rows=1, cols=1).cell(0, 0).text = "Nested"
        doc.save(docx_file)
        monkeypatch.setattr(document_utils, "_MAX_DOCX_TABLE_DEPTH", 1)

        result = extract_text(docx_file)

        assert result is not None
        assert result.startswith("[error: unsafe DOCX:")

    def test_extract_text_docx_empty(self, tmp_path: Path):
        """Test extracting text from an empty .docx file."""
        from docx import Document

        docx_file = tmp_path / "empty.docx"
        doc = Document()
        doc.save(docx_file)

        result = extract_text(docx_file)
        assert result == ""

    def test_extract_text_pptx(self, tmp_path: Path):
        """Test extracting text from a .pptx file."""
        from pptx import Presentation

        pptx_file = tmp_path / "test.pptx"
        prs = Presentation()

        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        for shape in slide1.shapes:
            if hasattr(shape, "text"):
                shape.text = "First Slide Title"

        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = width = height = 1000000
        textbox = slide2.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Bullet point content"

        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "--- Slide 1 ---" in result
        assert "--- Slide 2 ---" in result
        # Text content may vary depending on PowerPoint layout defaults
        assert len(result) > 0

    def test_extract_text_pptx_table(self, tmp_path: Path):
        """Table cells should be extracted, not silently dropped."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(1)
        ).table
        table.cell(0, 0).text = "Header A"
        table.cell(0, 1).text = "Header B"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "Bob"
        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "Header A" in result
        assert "Header B" in result
        assert "Alice" in result
        assert "Bob" in result

    def test_extract_text_pptx_grouped_shapes(self, tmp_path: Path):
        """Text inside grouped shapes must be extracted recursively."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "grouped.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        group = slide.shapes.add_group_shape()
        inner = group.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        inner.text_frame.text = "Inside group"
        prs.save(pptx_file)

        result = extract_text(pptx_file)
        assert result is not None
        assert "Inside group" in result

    def test_extract_text_rejects_oversized_office_archive(self, tmp_path, monkeypatch):
        office_file = tmp_path / "oversized.docx"
        with ZipFile(office_file, "w") as archive:
            archive.writestr("word/document.xml", "x" * 32)

        monkeypatch.setattr("nanobot.utils.document._MAX_OFFICE_UNCOMPRESSED_SIZE", 16)

        assert "Office document expands beyond" in (extract_text(office_file) or "")

    def test_extract_text_stops_streaming_xlsx_at_text_limit(self, tmp_path, monkeypatch):
        from openpyxl import Workbook, load_workbook

        xlsx_file = tmp_path / "large.xlsx"
        wb = Workbook(write_only=True)
        ws = wb.create_sheet()
        for index in range(100):
            ws.append([f"row-{index}-" + "x" * 20])
        wb.save(xlsx_file)

        visited = 0
        real_load_workbook = load_workbook

        def tracked_load_workbook(*args, **kwargs):
            workbook = real_load_workbook(*args, **kwargs)
            worksheet = workbook[workbook.sheetnames[0]]
            original_iter_rows = worksheet.iter_rows

            def tracked_rows(*row_args, **row_kwargs):
                nonlocal visited
                for row in original_iter_rows(*row_args, **row_kwargs):
                    visited += 1
                    yield row

            worksheet.iter_rows = tracked_rows
            return workbook

        monkeypatch.setattr("openpyxl.load_workbook", tracked_load_workbook)
        monkeypatch.setattr("nanobot.utils.document._MAX_TEXT_LENGTH", 80)

        result = extract_text(xlsx_file)

        assert result is not None
        assert "truncated at 80 chars" in result
        assert visited < 100

    def test_extract_pdf_pages_rejects_large_content_stream(self, tmp_path, monkeypatch):
        class _Contents:
            @staticmethod
            def get_data():
                return b"x" * 17

        class _Page:
            @staticmethod
            def get_contents():
                return _Contents()

            @staticmethod
            def extract_text():
                return "should not be reached"

        class _Reader:
            def __init__(self, *_args, **_kwargs):
                self.pages = [_Page()]

        monkeypatch.setattr("pypdf.PdfReader", _Reader)
        monkeypatch.setattr("nanobot.utils.document._MAX_PDF_CONTENT_STREAM_SIZE", 16)

        with pytest.raises(PdfSafetyError, match="content stream exceeds"):
            extract_pdf_pages(tmp_path / "large.pdf")

    def test_extract_text_pdf_not_found(self, tmp_path: Path):
        """Test that missing PDF files return error string."""
        missing_pdf = tmp_path / "nonexistent.pdf"

        result = extract_text(missing_pdf)
        assert result is not None
        assert "[error: file not found:" in result

    def test_extract_text_image_files(self, tmp_path: Path):
        """Test that image files return placeholder text."""
        # Create a minimal PNG file (1x1 pixel)
        png_file = tmp_path / "test.png"
        # Minimal valid PNG: 8-byte signature + IHDR + IDAT + IEND
        png_data = (
            b"\x89PNG\r\n\x1a\n"
            b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01"
            b"\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        png_file.write_bytes(png_data)

        result = extract_text(png_file)
        assert result is not None
        assert "[image:" in result
        assert "test.png" in result


class TestIsTextExtension:
    """Test the _is_text_extension helper."""

    def test_text_extensions_return_true(self):
        """Test that known text extensions return True."""
        assert _is_text_extension(".txt") is True
        assert _is_text_extension(".md") is True
        assert _is_text_extension(".csv") is True
        assert _is_text_extension(".json") is True
        assert _is_text_extension(".yaml") is True
        assert _is_text_extension(".yml") is True
        assert _is_text_extension(".xml") is True
        assert _is_text_extension(".html") is True
        assert _is_text_extension(".htm") is True

    def test_non_text_extensions_return_false(self):
        """Test that non-text extensions return False."""
        assert _is_text_extension(".pdf") is False
        assert _is_text_extension(".docx") is False
        assert _is_text_extension(".xlsx") is False
        assert _is_text_extension(".pptx") is False
        assert _is_text_extension(".png") is False
        assert _is_text_extension(".xyz") is False

    def test_case_sensitivity(self):
        """Test that _is_text_extension requires lowercase extension.

        Note: The main extract_text function handles case-insensitivity by
        converting extensions to lowercase before calling _is_text_extension.
        """
        # _is_text_extension itself is case-sensitive (lowercase only)
        assert _is_text_extension(".txt") is True
        assert _is_text_extension(".TXT") is False
        assert _is_text_extension(".pdf") is False
